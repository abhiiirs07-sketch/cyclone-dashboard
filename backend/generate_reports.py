import os
import csv
import json
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from pptx import Presentation
from pptx.util import Inches as PtInches, Pt as PtFont
from pptx.dml.color import RGBColor as PptRGBColor

# --- LOAD DATA ---
data = {
    "cyclone_name": "Fani",
    "landfall_place": "Puri",
    "rainfall_mean_mm": 54.0,
    "rainfall_max_mm": 162.0,
    "flooded_area_km2": 37.2,
    "vegetation_damaged_area_km2": 3207.7,
    "population_total": 27291030,
    "population_flooded": 4849,
    "top_districts": [
        {"name": "Khordha", "hazard": 0.198, "rain": 78.4, "flood": 8.5, "veg": 420.1, "pop": 2465180},
        {"name": "Cuttack", "hazard": 0.197, "rain": 69.2, "flood": 6.8, "veg": 380.4, "pop": 2624470},
        {"name": "Jajpur", "hazard": 0.196, "rain": 58.1, "flood": 5.4, "veg": 310.2, "pop": 1827192},
        {"name": "Baleshwar", "hazard": 0.194, "rain": 52.3, "flood": 4.1, "veg": 290.8, "pop": 2320500},
        {"name": "Bhadrak", "hazard": 0.193, "rain": 51.7, "flood": 3.9, "veg": 285.3, "pop": 1506300},
        {"name": "Kendrapara", "hazard": 0.193, "rain": 50.9, "flood": 3.7, "veg": 270.5, "pop": 1440360},
        {"name": "Jagatsinghpur", "hazard": 0.190, "rain": 49.5, "flood": 2.8, "veg": 250.6, "pop": 1136900},
        {"name": "Ganjam", "hazard": 0.186, "rain": 45.2, "flood": 1.2, "veg": 190.4, "pop": 3529031},
        {"name": "Dhenkanal", "hazard": 0.182, "rain": 42.1, "flood": 0.5, "veg": 140.2, "pop": 1192700},
        {"name": "Srikakulam", "hazard": 0.182, "rain": 41.8, "flood": 0.3, "veg": 135.1, "pop": 2704700}
    ]
}

# --- GENERATE CSV ---
csv_path = "cyclone_fani_chart_data.csv"
print(f"Generating {csv_path}...")
with open(csv_path, mode="w", newline="", encoding="utf-8") as f:
    writer = csv.writer(f)
    writer.writerow(["District", "Composite Hazard Index", "Mean Rainfall (mm)", "Flooded Area (km2)", "Vegetation Damage (km2)", "Population Exposed"])
    for d in data["top_districts"]:
        writer.writerow([d["name"], d["hazard"], d["rain"], d["flood"], d["veg"], d["pop"]])

# --- GENERATE WORD DOCUMENT ---
docx_path = "cyclone_fani_impact_report.docx"
print(f"Generating {docx_path}...")
doc = Document()

# Styles & Colors
title_style = doc.styles['Title']
title_style.font.name = 'Arial'
title_style.font.size = Pt(26)
title_style.font.bold = True
title_style.font.color.rgb = RGBColor(15, 23, 42)

doc.add_heading("Cyclone Fani Impact Assessment Report", level=0)
doc.add_paragraph("Google Earth Engine Powered Disaster Monitoring System\nGenerated: 2026")

# Executive Summary
doc.add_heading("1. Executive Summary", level=1)
doc.add_paragraph(
    "This report provides an automated geospatial impact assessment for Cyclone Fani (2019) "
    "across coastal regions of Odisha and Andhra Pradesh, India. Built using the Google Earth Engine "
    "API and multi-satellite sensors, the platform maps storm inundation, vegetation damage, "
    "meteorological severity, and population exposure."
)

# Key Stats Table
table = doc.add_table(rows=1, cols=2)
table.style = 'Light Shading Accent 1'
hdr_cells = table.rows[0].cells
hdr_cells[0].text = 'Indicator / Metric'
hdr_cells[1].text = 'Value / Extent'

metrics = [
    ("Landfall Location", data["landfall_place"]),
    ("Mean Precipitation (GPM)", f"{data['rainfall_mean_mm']} mm"),
    ("Max Precipitation (GPM)", f"{data['rainfall_max_mm']} mm"),
    ("Flooded Inundation Area (Sentinel-1 SAR)", f"{data['flooded_area_km2']} km²"),
    ("Vegetation Damage Area (NDVI)", f"{data['vegetation_damaged_area_km2']} km²"),
    ("Total Exposed Population", f"{data['population_total']:,} people"),
    ("Directly Inundated Population", f"{data['population_flooded']:,} people")
]

for metric, val in metrics:
    row_cells = table.add_row().cells
    row_cells[0].text = metric
    row_cells[1].text = val

doc.add_heading("2. District-Level Exposure Breakdown", level=1)
doc.add_paragraph("The table below details the top affected districts ranked by their Composite Hazard Index:")

dist_table = doc.add_table(rows=1, cols=6)
dist_table.style = 'Table Grid'
hdr = dist_table.rows[0].cells
hdr[0].text = 'District'
hdr[1].text = 'Hazard Index'
hdr[2].text = 'Rainfall (mm)'
hdr[3].text = 'Flood (km²)'
hdr[4].text = 'Veg Damage (km²)'
hdr[5].text = 'Exposed Pop'

for d in data["top_districts"]:
    row = dist_table.add_row().cells
    row[0].text = d["name"]
    row[1].text = f"{d['hazard']:.3f}"
    row[2].text = f"{d['rain']} mm"
    row[3].text = f"{d['flood']} km²"
    row[4].text = f"{d['veg']} km²"
    row[5].text = f"{d['pop']:,}"

doc.save(docx_path)

# --- GENERATE POWERPOINT PRESENTATION ---
pptx_path = "cyclone_fani_presentation.pptx"
print(f"Generating {pptx_path}...")
prs = Presentation()

def add_slide_layout(title_text, content_items, icon="🌪️"):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Gradient header block representation
    left = PtInches(0.5)
    top = PtInches(0.5)
    width = PtInches(9.0)
    height = PtInches(1.0)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"{icon} {title_text}"
    p.font.size = PtFont(32)
    p.font.bold = True
    p.font.color.rgb = PptRGBColor(15, 23, 42)
    
    # Body
    left_body = PtInches(0.5)
    top_body = PtInches(1.8)
    width_body = PtInches(9.0)
    height_body = PtInches(5.0)
    bodyBox = slide.shapes.add_textbox(left_body, top_body, width_body, height_body)
    tf_body = bodyBox.text_frame
    tf_body.word_wrap = True
    
    for i, item in enumerate(content_items):
        p_item = tf_body.add_paragraph() if i > 0 else tf_body.paragraphs[0]
        p_item.text = f"• {item}"
        p_item.font.size = PtFont(16)
        p_item.font.color.rgb = PptRGBColor(75, 85, 99)
        p_item.space_after = PtFont(14)

# Slide 1: Title
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)
title_box = slide.shapes.add_textbox(PtInches(0.5), PtInches(2.5), PtInches(9.0), PtInches(3.0))
tf = title_box.text_frame
p1 = tf.paragraphs[0]
p1.text = "Cyclone Fani Disaster Impact Report"
p1.font.size = PtFont(40)
p1.font.bold = True
p1.font.color.rgb = PptRGBColor(37, 99, 235)

p2 = tf.add_paragraph()
p2.text = "Satellite-Based Environmental and Infrastructure Assessment Dashboard"
p2.font.size = PtFont(18)
p2.font.color.rgb = PptRGBColor(107, 114, 128)
p2.space_before = PtFont(10)

# Slide 2: Problem
add_slide_layout("The Cyclone Disaster Challenge", [
    "Cyclone Fani hit Odisha in May 2019 with landfall at Puri, causing severe destruction.",
    " NDRF and relief workers lacked unified, real-time spatial analysis of flood extent and infrastructure damage.",
    "Manual satellite data processing takes days — losing the critical 72-hour window for life-saving operations."
], "❌")

# Slide 3: GEE Analytics
add_slide_layout("Google Earth Engine Core Engine", [
    "Calculates pre/post SAR difference dynamically using Sentinel-1 C-band radar backscatter.",
    "Automatically thresholds water surfaces to define absolute flooded boundaries (37.2 km² mapped).",
    "Runs multi-criteria overlay model combining wind buffers, inundation zones, and population exposure."
], "🛰️")

# Slide 4: Results Overview
add_slide_layout("Key Assessment Results", [
    "Rainfall: GPM recorded up to 162.0 mm of precipitation during landfall.",
    "Inundated Land: Mapped 37.2 km² of flooded regions via SAR sensor analysis.",
    "Vegetation Damage: NDVI degradation mapping identified 3207.7 km² of damaged agricultural land.",
    "Population Impacted: Mapped 27+ million exposed coastal residents (with 4,849 in high-risk flood zones)."
], "📈")

# Slide 5: Top Districts Table
add_slide_layout("High-Exposure Districts", [
    "Khordha: Ranked #1 with a Hazard Index of 0.198 (8.5 km² flood, 2.4M exposed population).",
    "Cuttack: Ranked #2 with a Hazard Index of 0.197 (6.8 km² flood, 2.6M exposed population).",
    "Jajpur: Ranked #3 with a Hazard Index of 0.196 (5.4 km² flood, 1.8M exposed population).",
    "Baleshwar & Bhadrak: Severe storm surge damage corridors mapped along coastal zones."
], "🗺️")

# Slide 6: Tech stack
add_slide_layout("Full Cloud Infrastructure", [
    "Frontend: Next.js 15 statically compiled and hosted on Netlify CDN for maximum speed.",
    "Backend: FastAPI async server hosted on Railway container engine.",
    "Earth Engine Integration: Authenticated via OAuth refresh token credentials.",
    "Dynamic API: Serves XYZ satellite map tile URLs directly to MapLibre GL UI."
], "💻")

prs.save(pptx_path)
print("All reports generated successfully!")
